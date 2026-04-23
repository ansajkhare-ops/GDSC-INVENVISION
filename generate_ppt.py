# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color Palette ──────────────────────────────────────────
BG      = RGBColor(0x0D, 0x10, 0x17)   # dark navy
ACCENT  = RGBColor(0x00, 0xD4, 0xFF)   # cyan
GREEN   = RGBColor(0x00, 0xFF, 0x9D)   # green
ORANGE  = RGBColor(0xFF, 0x6B, 0x35)   # orange
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xB0, 0xBE, 0xC5)
RED     = RGBColor(0xFF, 0x3D, 0x3D)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)

BLANK = prs.slide_layouts[6]  # completely blank


def slide():
    return prs.slides.add_slide(BLANK)


def bg(sl, color=BG):
    sh = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def box(sl, x, y, w, h, color, alpha=None):
    sh = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(sl, text, x, y, w, h, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color
    return tb


def accent_bar(sl, y=0.55, color=ACCENT):
    box(sl, 0.7, y, 1.1, 0.07, color)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
# Gradient top strip
box(s, 0, 0, 13.33, 0.08, ACCENT)
# Big title
txt(s, "InvenVision", 1.5, 1.5, 10, 1.8, size=72, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "AI-Powered Inventory Intelligence for Every Business", 1.5, 3.1, 10, 0.8,
    size=26, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Predict  ·  Manage  ·  Never Run Out", 1.5, 3.85, 10, 0.6,
    size=18, color=LGRAY, align=PP_ALIGN.CENTER)
# Bottom bar
box(s, 0, 7.1, 13.33, 0.4, RGBColor(0x06, 0x08, 0x0F))
txt(s, "GDSC Hackathon 2026  |  Team InvenVision  |  Live at gdsc-hackathon.up.railway.app",
    0.5, 7.1, 12.33, 0.4, size=13, color=LGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — What is Supply Chain?
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, GREEN)
accent_bar(s, color=GREEN)
txt(s, "Understanding the Supply Chain", 0.7, 0.6, 12, 0.7, size=34,
    bold=True, color=WHITE)
txt(s, "The journey of a product — from manufacturer to your hands",
    0.7, 1.25, 12, 0.5, size=17, color=LGRAY)

# Flow boxes
steps = [("🏭", "Manufacturer", "Produces goods\nin bulk"),
         ("🚚", "Distributor", "Transports & stores\nin warehouses"),
         ("🏪", "Retailer", "Sells to end\ncustomers"),
         ("👤", "Customer", "Receives the\nproduct")]
colors = [ACCENT, GREEN, ORANGE, YELLOW]
for i, (icon, title, desc) in enumerate(steps):
    x = 0.7 + i * 3.1
    box(s, x, 2.1, 2.8, 2.5, RGBColor(0x16, 0x1B, 0x22))
    txt(s, icon,  x+0.9, 2.2, 1.2, 0.6, size=32, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.1, 2.85, 2.6, 0.5, size=17, bold=True,
        color=colors[i], align=PP_ALIGN.CENTER)
    txt(s, desc,  x+0.1, 3.35, 2.6, 0.9, size=13, color=LGRAY,
        align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "→", x+2.83, 2.9, 0.4, 0.5, size=28, color=ACCENT,
            align=PP_ALIGN.CENTER)

txt(s, "Applies to every business — Pharmacy, Grocery, Retail, Electronics, Restaurant & more.",
    0.7, 5.0, 12, 0.6, size=16, color=LGRAY)
txt(s, "Break one link → the entire chain fails.", 0.7, 5.55, 12, 0.5,
    size=18, bold=True, color=ORANGE)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — Problems in Supply Chain
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, RED)
accent_bar(s, color=RED)
txt(s, "Critical Problems in Supply Chain Management",
    0.7, 0.6, 12, 0.7, size=32, bold=True, color=WHITE)

problems = [
    ("📉", "Demand Forecasting Failure",
     "Businesses can't predict demand → overstock or stockout — affects every sector"),
    ("🗑️", "Waste & Dead Stock",
     "Perishables, medicines, seasonal goods expire → ₹15,000+ loss per shop per year"),
    ("🔗", "No Real-Time Visibility",
     "No live view of stock levels — owners fly blind across all product categories"),
    ("📝", "Manual Tracking",
     "Paper ledgers & Excel → errors, delays, data loss in grocery, retail & pharmacy alike"),
    ("💸", "Cash Flow Locked in Dead Stock",
     "Overstocked products tie up capital that every small business desperately needs"),
    ("⏰", "Reactive Decisions Only",
     "Owners in every industry reorder only after running out — always one step behind"),
]
for i, (icon, title, desc) in enumerate(problems):
    row, col = divmod(i, 2)
    x = 0.6 + col * 6.3
    y = 1.6 + row * 1.6
    box(s, x, y, 5.9, 1.35, RGBColor(0x20, 0x08, 0x08))
    txt(s, icon + "  " + title, x+0.2, y+0.1, 5.5, 0.5,
        size=16, bold=True, color=RED)
    txt(s, desc, x+0.2, y+0.6, 5.5, 0.6, size=13, color=LGRAY)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — The Pharmacy-Specific Problem
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, ORANGE)
accent_bar(s, color=ORANGE)
txt(s, "The Small Business Owner's Daily Nightmare",
    0.7, 0.6, 12, 0.7, size=34, bold=True, color=WHITE)

# Industry tags
industries = [("💊", "Pharmacy"), ("🛒", "Grocery"), ("👗", "Retail"),
              ("📱", "Electronics"), ("🍽️", "Restaurant"), ("📦", "Wholesale")]
for i, (ic, nm) in enumerate(industries):
    x = 0.55 + i * 2.15
    box(s, x, 1.35, 1.9, 0.55, RGBColor(0x22, 0x14, 0x04))
    txt(s, ic + " " + nm, x+0.1, 1.43, 1.72, 0.38, size=13,
        bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Stat cards
stats = [("63 Lakh+", "MSMEs in India\nwith inventory problems"),
         ("87%", "Still using paper\nor basic Excel"),
         ("₹20,000", "Average annual loss\nto dead/expired stock"),
         ("1 in 3", "Customers lost due\nto stockouts")]
for i, (num, label) in enumerate(stats):
    x = 0.55 + i * 3.1
    box(s, x, 2.2, 2.85, 1.85, RGBColor(0x1A, 0x10, 0x04))
    txt(s, num,   x+0.15, 2.35, 2.55, 0.9, size=36, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, label, x+0.15, 3.2,  2.55, 0.7, size=12, color=LGRAY,
        align=PP_ALIGN.CENTER)

box(s, 0.6, 4.35, 12.1, 1.2, RGBColor(0x1A, 0x10, 0x04))
txt(s, "😫  The reality:", 0.9, 4.45, 3, 0.5, size=17, bold=True, color=ORANGE)
txt(s, '"I ordered too much stock last month — half of it expired or became dead stock.\n'
       'This month I ran out of my bestseller in 2 days and lost 30 customers."',
    3.5, 4.45, 9, 0.9, size=14, color=LGRAY)
txt(s, "This is not just a pharmacy problem. This is every small business owner in India. Every month.",
    0.7, 5.75, 12, 0.5, size=16, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — OUR SOLUTION
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, ACCENT)
accent_bar(s, color=ACCENT)
txt(s, "Introducing InvenVision", 0.7, 0.6, 12, 0.7, size=38,
    bold=True, color=WHITE)
txt(s, "One platform. Every business. Zero stock surprises.", 0.7, 1.25, 12, 0.5,
    size=20, color=ACCENT)

features = [
    ("🔮", "AI Demand Forecasting",
     "6 models predict when ANY product will run out — pharmacy, grocery, retail & more"),
    ("🚨", "Plain-English Alerts",
     '"Your Rice stock will last only 2 more days. Order NOW." — works for any product'),
    ("📊", "Smart Analytics",
     "Revenue trends, top products, category insights — from real sales data"),
    ("🧾", "Full Billing System",
     "Create invoices, auto-update inventory, track payments — for any category"),
    ("📒", "Customer Khata",
     "Digital credit ledger — works for every business that gives credit to customers"),
    ("☁️", "Cloud-First & Multi-User",
     "Any device, anywhere — one account per shop, Supabase PostgreSQL, zero data loss"),
]
for i, (icon, title, desc) in enumerate(features):
    row, col = divmod(i, 2)
    x = 0.6 + col * 6.3
    y = 2.05 + row * 1.55
    box(s, x, y, 5.9, 1.35, RGBColor(0x04, 0x15, 0x22))
    txt(s, icon + "  " + title, x+0.2, y+0.1, 5.5, 0.5,
        size=16, bold=True, color=ACCENT)
    txt(s, desc, x+0.2, y+0.58, 5.5, 0.65, size=13, color=LGRAY)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — AI FORECASTING DEEP DIVE
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, GREEN)
accent_bar(s, color=GREEN)
txt(s, "AI Forecasting Engine — How It Works",
    0.7, 0.6, 12, 0.7, size=34, bold=True, color=WHITE)

models = [
    ("EMA", "Exponential Moving Average", "Weights recent sales more — best for trending products"),
    ("SMA", "Simple Moving Average",      "Smooth baseline for stable products"),
    ("WMA", "Weighted Moving Average",    "Custom weights for seasonal patterns"),
    ("Holt's", "Double Exponential",      "Captures both level & trend simultaneously"),
    ("Linear", "Linear Regression",       "Best for steady growth or decline"),
    ("Seasonal", "Naive Seasonal",        "Repeating weekly/monthly cycles"),
]
for i, (short, full, desc) in enumerate(models):
    row, col = divmod(i, 2)
    x = 0.55 + col * 6.35
    y = 1.6 + row * 1.5
    box(s, x, y, 0.9, 1.2, RGBColor(0x00, 0x25, 0x1A))
    txt(s, short, x+0.05, y+0.28, 0.82, 0.65, size=13, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, full,  x+1.0,  y+0.05, 4.8, 0.5, size=15, bold=True, color=WHITE)
    txt(s, desc,  x+1.0,  y+0.55, 4.8, 0.55, size=13, color=LGRAY)

box(s, 0.55, 6.25, 12.2, 0.7, RGBColor(0x04, 0x20, 0x12))
txt(s, "✅  System auto-selects the best model per product. User just picks a product → clicks Generate → sees results in seconds.",
    0.8, 6.32, 11.8, 0.55, size=14, color=GREEN)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM FLOW
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, ACCENT)
accent_bar(s, color=ACCENT)
txt(s, "System Architecture & Data Flow",
    0.7, 0.6, 12, 0.7, size=34, bold=True, color=WHITE)

# Flow
flow = [
    ("👤\nUser Login", "Google OAuth\nSecure Auth"),
    ("🧾\nRecord Sale", "Invoice Created\nStock Updated"),
    ("🗄️\nSupabase DB", "PostgreSQL\nCloud Storage"),
    ("🧠\nAI Engine", "6 Models Run\nBest Selected"),
    ("📊\nDashboard", "Alerts + Charts\n+ Forecast"),
]
fc = [ACCENT, GREEN, YELLOW, ORANGE, RGBColor(0xC0, 0x00, 0xFF)]
for i, (top, bot) in enumerate(flow):
    x = 0.5 + i * 2.55
    box(s, x, 1.8, 2.2, 2.0, RGBColor(0x0A, 0x15, 0x22))
    txt(s, top, x+0.1, 1.9, 2.0, 1.0, size=14, bold=True,
        color=fc[i], align=PP_ALIGN.CENTER)
    txt(s, bot, x+0.1, 2.9, 2.0, 0.7, size=12, color=LGRAY,
        align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", x+2.22, 2.45, 0.4, 0.5, size=24, color=ACCENT,
            align=PP_ALIGN.CENTER)

# DB explanation
txt(s, "Database Tables:", 0.7, 4.2, 3, 0.45, size=16, bold=True, color=ACCENT)
tables = "users  ·  products  ·  stock_in  ·  invoices  ·  invoice_items  ·  customers  ·  customer_payments  ·  inventory"
txt(s, tables, 0.7, 4.65, 12, 0.5, size=14, color=LGRAY)
txt(s, "All tables auto-created on first boot. Zero manual DB setup required.",
    0.7, 5.1, 12, 0.4, size=13, color=GREEN)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, YELLOW)
accent_bar(s, color=YELLOW)
txt(s, "Technology Stack", 0.7, 0.6, 12, 0.7, size=38, bold=True, color=WHITE)

layers = [
    ("🖥️  Frontend", ACCENT,
     ["HTML5 + CSS3 (Vanilla)", "JavaScript (ES6+)", "Chart.js — interactive graphs",
      "Responsive design — mobile ready"]),
    ("⚙️  Backend", GREEN,
     ["Python 3.12 + Flask", "Authlib — Google OAuth 2.0", 
      "Gunicorn — production WSGI", "ProxyFix — HTTPS behind Railway"]),
    ("🗄️  Database", YELLOW,
     ["PostgreSQL (Supabase cloud)", "psycopg2-binary driver",
      "Session Pooler — IPv4, SSL", "8 relational tables"]),
    ("☁️  DevOps", ORANGE,
     ["Railway — cloud deployment", "GitHub — CI/CD auto-deploy",
      "nixpacks.toml — build config", "Environment vars — zero secrets in code"]),
]
for i, (title, color, items) in enumerate(layers):
    col = i % 2; row = i // 2
    x = 0.55 + col * 6.35
    y = 1.55 + row * 2.65
    box(s, x, y, 6.1, 2.45, RGBColor(0x0A, 0x12, 0x1A))
    box(s, x, y, 6.1, 0.5, RGBColor(0x0F, 0x1E, 0x2E))
    txt(s, title, x+0.2, y+0.06, 5.7, 0.42, size=17, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s, "▸  " + item, x+0.2, y+0.6+j*0.44, 5.7, 0.42, size=13, color=LGRAY)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES SHOWCASE
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, RGBColor(0xC0, 0x00, 0xFF))
accent_bar(s, color=RGBColor(0xC0, 0x00, 0xFF))
txt(s, "What Makes InvenVision Different",
    0.7, 0.6, 12, 0.7, size=34, bold=True, color=WHITE)

compare = [
    ("Feature", "Traditional Tools", "InvenVision", True),
    ("Demand Forecasting", "❌  None", "✅  6 AI Models", False),
    ("Plain-English Alerts", "❌  Only numbers/reports", "✅  'Order NOW — 2 days left'", False),
    ("Cloud Database", "❌  Local PC only", "✅  Supabase PostgreSQL", False),
    ("Google Login", "❌  Username+Password", "✅  One-click OAuth", False),
    ("Analytics Charts", "❌  Export to Excel only", "✅  Live interactive charts", False),
    ("Customer Khata", "❌  Separate notebook", "✅  Digital credit ledger built-in", False),
]
for i, (feat, old, new, header) in enumerate(compare):
    y = 1.45 + i * 0.65
    rc = RGBColor(0x12, 0x18, 0x26) if i % 2 == 0 else RGBColor(0x0D, 0x12, 0x1E)
    if header:
        rc = RGBColor(0x1A, 0x06, 0x2E)
    box(s, 0.5, y, 4.0, 0.6, rc)
    box(s, 4.55, y, 3.8, 0.6, rc)
    box(s, 8.4, y, 4.4, 0.6, rc)
    c = RGBColor(0xC0, 0x00, 0xFF) if header else LGRAY
    cn = GREEN if not header else RGBColor(0xC0, 0x00, 0xFF)
    txt(s, feat, 0.6, y+0.1, 3.8, 0.45, size=13, bold=header, color=c)
    txt(s, old,  4.65, y+0.1, 3.6, 0.45, size=13, bold=header, color=RED if not header else c)
    txt(s, new,  8.5, y+0.1, 4.2, 0.45, size=13, bold=header, color=cn)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — LIVE DEPLOYMENT
# ══════════════════════════════════════════════════════════

s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, ACCENT)
accent_bar(s, color=ACCENT)
txt(s, "Live & Deployed — Right Now",
    0.7, 0.6, 12, 0.7, size=38, bold=True, color=WHITE)

urls = [
    ("🌐", "Live App",      "gdsc-hackathon.up.railway.app", ACCENT),
    ("💚", "Health Check",  "gdsc-hackathon.up.railway.app/api/health", GREEN),
    ("💻", "Source Code",   "github.com/ansajkhare-ops/GDSC-INVENVISION", YELLOW),
    ("🗄️", "Database",      "Supabase — aws-1-ap-southeast-1 (PostgreSQL)", ORANGE),
]
for i, (icon, label, url, color) in enumerate(urls):
    y = 1.8 + i * 1.1
    box(s, 0.55, y, 12.2, 0.9, RGBColor(0x04, 0x14, 0x22))
    txt(s, icon + "  " + label, 0.8, y+0.15, 3.0, 0.6,
        size=16, bold=True, color=color)
    txt(s, url, 3.9, y+0.2, 8.7, 0.5, size=15, color=WHITE)

box(s, 0.55, 6.35, 12.2, 0.65, RGBColor(0x00, 0x25, 0x10))
txt(s, "✅  Not a prototype. Not a mockup. A real, production-deployed application.",
    0.8, 6.42, 11.8, 0.5, size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING
# ══════════════════════════════════════════════════════════
s = slide(); bg(s)
box(s, 0, 0, 13.33, 0.08, ACCENT)
box(s, 0, 7.42, 13.33, 0.08, ACCENT)

txt(s, "InvenVision", 1.5, 1.3, 10, 1.4, size=68, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "The Brain Every Business is Missing",
    1.5, 2.7, 10, 0.7, size=26, color=ACCENT, align=PP_ALIGN.CENTER)

box(s, 3.5, 3.65, 6.3, 0.06, RGBColor(0x1A, 0x2A, 0x3A))

quotes = [
    "🔮  Predict demand before stock runs out — any product, any category",
    "🚨  Get plain-English alerts — not numbers, not reports",
    "📊  See real analytics from real sales data — automatically",
    "☁️  One platform for Pharmacy, Grocery, Retail, Electronics & more",
]
for i, q in enumerate(quotes):
    txt(s, q, 2.0, 4.0 + i*0.55, 9.3, 0.5, size=15, color=LGRAY,
        align=PP_ALIGN.CENTER)

txt(s, "Thank You", 1.5, 6.5, 10, 0.8, size=32, bold=True,
    color=ACCENT, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────
out = r"C:\Users\DELL\Invenvision\InvenVision_Pitch.pptx"
prs.save(out)
print("PPT SAVED:", out)
